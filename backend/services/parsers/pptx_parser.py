from io import BytesIO
from pptx import Presentation

from .file_parser import FileParser

class PptxParser(FileParser):

    async def parse_text(self) -> str:
        content = await self.file.read()
        prs = Presentation(BytesIO(content))

        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)

        return " ".join(text_runs)